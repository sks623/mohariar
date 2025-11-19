"""
MOHARIAR PowerPoint Presentation Generator
Creates professional presentation/user guide
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Define colors (MOHARIAR theme)
DARK_BG = RGBColor(10, 14, 39)
IRON_RED = RGBColor(185, 29, 29)
IRON_GOLD = RGBColor(255, 215, 0)
STARK_BLUE = RGBColor(0, 168, 232)
WHITE = RGBColor(255, 255, 255)

def add_title_slide(prs, title, subtitle):
    """Add title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank

    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BG

    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_p = title_frame.paragraphs[0]
    title_p.alignment = PP_ALIGN.CENTER
    title_p.font.size = Pt(48)
    title_p.font.bold = True
    title_p.font.color.rgb = IRON_GOLD

    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(8), Inches(1))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = subtitle
    subtitle_p = subtitle_frame.paragraphs[0]
    subtitle_p.alignment = PP_ALIGN.CENTER
    subtitle_p.font.size = Pt(24)
    subtitle_p.font.color.rgb = WHITE

def add_content_slide(prs, title, content_lines):
    """Add content slide with bullet points"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BG

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_p = title_frame.paragraphs[0]
    title_p.font.size = Pt(32)
    title_p.font.bold = True
    title_p.font.color.rgb = IRON_GOLD

    # Content
    content_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.5), Inches(8.5), Inches(5.5))
    content_frame = content_box.text_frame
    content_frame.word_wrap = True

    for i, line in enumerate(content_lines):
        if i > 0:
            content_frame.add_paragraph()
        p = content_frame.paragraphs[i]
        p.text = line
        p.font.size = Pt(16)
        p.font.color.rgb = WHITE
        p.space_after = Pt(8)

        # Indent sub-bullets
        if line.startswith('  •') or line.startswith('  ├') or line.startswith('  └'):
            p.level = 1
            p.font.size = Pt(14)

def add_screenshot_slide(prs, title, screenshot_description, content_lines=None):
    """Add slide with screenshot placeholder"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BG

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_p = title_frame.paragraphs[0]
    title_p.font.size = Pt(28)
    title_p.font.bold = True
    title_p.font.color.rgb = IRON_GOLD

    # Screenshot placeholder
    screenshot_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(3))
    screenshot_frame = screenshot_box.text_frame
    screenshot_frame.text = f"[INSERT SCREENSHOT: {screenshot_description}]"
    screenshot_p = screenshot_frame.paragraphs[0]
    screenshot_p.alignment = PP_ALIGN.CENTER
    screenshot_p.font.size = Pt(18)
    screenshot_p.font.italic = True
    screenshot_p.font.color.rgb = STARK_BLUE

    # Add border to screenshot area
    shape = slide.shapes.add_shape(
        1, Inches(1), Inches(1.5), Inches(8), Inches(3)  # Rectangle
    )
    shape.fill.background()
    shape.line.color.rgb = STARK_BLUE
    shape.line.width = Pt(2)

    # Content below screenshot
    if content_lines:
        content_box = slide.shapes.add_textbox(Inches(0.75), Inches(5), Inches(8.5), Inches(2))
        content_frame = content_box.text_frame
        content_frame.word_wrap = True

        for i, line in enumerate(content_lines):
            if i > 0:
                content_frame.add_paragraph()
            p = content_frame.paragraphs[i]
            p.text = line
            p.font.size = Pt(14)
            p.font.color.rgb = WHITE

# SLIDE 1: Title
add_title_slide(
    prs,
    "⚡ MOHARIAR ⚡",
    "RI Online Collection Automation System\n\nDesigned by: SUSHANT\nVersion: 2.0 | November 2025"
)

# SLIDE 2: What is MOHARIAR?
add_content_slide(prs, "What is MOHARIAR?", [
    "🎯 PURPOSE",
    "Automates bulk processing of land revenue collection entries",
    "on Odisha Land Revenue Portal",
    "",
    "✅ KEY BENEFITS",
    "• Process 100+ khatas without manual intervention",
    "• Auto-save PDFs with organized folder structure",
    "• Handle multiple villages simultaneously",
    "• Real-time progress tracking with color-coded statistics",
    "• 95%+ success rate with smart retry mechanism",
    "",
    "⏱️ TIME SAVINGS",
    "• Manual: 3-5 minutes per khata",
    "• MOHARIAR: 8-12 seconds per khata",
    "• Result: 20x FASTER!"
])

# SLIDE 3: System Requirements
add_content_slide(prs, "System Requirements", [
    "💻 SOFTWARE REQUIREMENTS",
    "• Windows 7 / 10 / 11",
    "• Google Chrome Browser (latest version)",
    "• Odisha Land Revenue Portal credentials",
    "• Stable internet connection",
    "",
    "📦 WHAT'S INCLUDED",
    "• MOHARIAR.exe (standalone executable - no Python needed)",
    "• This User Guide (PDF format)",
    "• Welcome & Completion sound notifications",
    "• Excel template for bulk upload",
    "• Auto-organized PDF session folders"
])

# SLIDE 4: Key Features
add_content_slide(prs, "Key Features", [
    "🌟 CORE AUTOMATION FEATURES",
    "✅ WebDriverWait Technology - Smart adaptive delays (40-50% faster)",
    "✅ CDP PDF Generation - Direct browser PDF creation (95%+ reliable)",
    "✅ Session Folders - Auto-organized by date and village",
    "✅ Excel Integration - Bulk upload khata lists with mobile numbers",
    "✅ Multiple Credentials - Save and switch between user profiles",
    "✅ Session History - Quick access to last 10 processing sessions",
    "",
    "📊 VISUAL & USER EXPERIENCE",
    "🟢 Color-coded live statistics (Success/Already Paid/Errors/PDFs)",
    "🎯 Real-time progress tracking with ETA countdown",
    "📄 Automatic PDF merge (per village + master file)",
    "🔔 Windows toast notifications on completion",
    "🎵 Sound alerts for important events"
])

# SLIDE 5: UI Overview (with screenshot)
add_screenshot_slide(
    prs,
    "User Interface Overview",
    "Full application window showing 3-panel layout",
    [
        "THREE-PANEL DESIGN:",
        "├─ LEFT (30%): Village Selection with Select All / Clear All",
        "├─ MIDDLE (40%): Khata Selection with search & per-village controls",
        "└─ RIGHT (30%): Live Stats, Progress Bar, and Mission Log"
    ]
)

# SLIDE 6: Getting Started
add_content_slide(prs, "Getting Started (First Time Setup)", [
    "STEP 1: Launch MOHARIAR.exe",
    "  • Double-click the executable file",
    "",
    "STEP 2: Enter Default Mobile Number",
    "  • Enter 10-digit mobile number (used as default)",
    "",
    "STEP 3: Click '🔐 Initialize System'",
    "  • Chrome browser will open automatically",
    "",
    "STEP 4: Login Manually (one-time setup)",
    "  • Select your District from dropdown",
    "  • Enter User ID and Password",
    "  • Solve the captcha puzzle",
    "  • Click Login button",
    "",
    "STEP 5: System auto-navigates to RI Collection module",
    "",
    "STEP 6: Ready to process! ✅"
])

# SLIDE 7: Basic Workflow
add_content_slide(prs, "Basic Workflow (Standard Processing)", [
    "1️⃣ SELECT VILLAGES",
    "  • Check boxes next to desired villages, OR",
    "  • Click 'Select All Villages' button",
    "",
    "2️⃣ EXTRACT KHATAS",
    "  • Click '📥 Extract Khatas' button",
    "  • System scans each selected village",
    "  • All available khatas are displayed",
    "",
    "3️⃣ SELECT KHATAS",
    "  • Use search box to filter (optional)",
    "  • Check individual khatas, OR",
    "  • Use 'Select All' button per village",
    "",
    "4️⃣ ENGAGE PROTOCOL",
    "  • Click '⚡ Engage Protocol' button",
    "  • Watch real-time progress with color-coded stats",
    "  • PDFs auto-save to session folder",
    "",
    "5️⃣ COMPLETION",
    "  • View final statistics",
    "  • Click 'Open Folder' to see PDFs",
    "  • Export results to Excel if needed"
])

# SLIDE 8: Excel Bulk Upload
add_content_slide(prs, "Excel Bulk Upload Feature", [
    "📥 UPLOAD KHATA LISTS VIA EXCEL",
    "",
    "STEP 1: Download Template",
    "  • Click 'Download Excel Template' button",
    "  • Opens: Khata_Template.xlsx",
    "",
    "STEP 2: Fill Template",
    "  • Column A: Khata Number (required)",
    "  • Column B: Mobile Number (optional - 10 digits or blank)",
    "",
    "STEP 3: Upload Per Village",
    "  • Check the village checkbox",
    "  • Click 'Upload Excel' button next to that village",
    "  • Select your filled Excel file",
    "  • Repeat for other villages if needed",
    "",
    "STEP 4: Process",
    "  • Click 'Extract Khatas' to verify data",
    "  • Click 'Engage Protocol' to process",
    "",
    "✅ Supports: .xlsx and .csv formats",
    "⚠️ Mobile: 10 digits or leave blank (uses default)"
])

# SLIDE 9: PDF Management
add_content_slide(prs, "PDF Management & Organization", [
    "📁 SESSION FOLDER STRUCTURE:",
    "",
    "rr20251117_1430/                    ← Session ID (date_time)",
    "├── VillageName1/",
    "│   ├── VillageName1_101.pdf",
    "│   ├── VillageName1_102.pdf",
    "│   └── VillageName1_14_23.pdf      ← Slash sanitized (14/23 → 14_23)",
    "├── VillageName2/",
    "│   └── VillageName2_201.pdf",
    "├── VillageName1_MERGED.pdf         ← Per-village merged PDF",
    "└── ALL_RECEIPTS_MERGED.pdf         ← Master merged PDF",
    "",
    "📊 PDF FEATURES:",
    "✅ Auto-save during processing (CDP technology - no dialogs)",
    "✅ Organized by village in subfolders",
    "✅ Auto-merge at completion (per village + master)",
    "✅ 'Open Folder' button for quick access",
    "✅ Manual 'Merge PDFs' available anytime"
])

# SLIDE 10: Live Statistics
add_screenshot_slide(
    prs,
    "Live Statistics (Color-Coded)",
    "Right panel showing color-coded statistics",
    [
        "📊 REAL-TIME STATUS DISPLAY:",
        "🟢 Success: X         (Green - successfully processed)",
        "🟡 Already Paid: X    (Yellow - skipped, no retry)",
        "🔴 Errors: X          (Red - failed after 3 retries)",
        "📄 PDFs: X/Y          (Blue - saved/total)",
        "",
        "Progress Bar: Shows percentage complete",
        "⏰ ETA: Estimated time remaining (updates live)",
        "Current Status: Shows which khata is processing"
    ]
)

# SLIDE 11: Multiple Credentials
add_content_slide(prs, "Multiple Credentials Feature", [
    "👥 PROFILE MANAGEMENT",
    "",
    "SAVE MULTIPLE USER PROFILES:",
    "• Profile 1: User_RTO",
    "• Profile 2: User_Backup",
    "• Profile 3: User_District2",
    "",
    "FEATURES:",
    "✅ Encrypted storage (secure credentials)",
    "✅ Quick profile switching via dropdown",
    "✅ Add / Edit / Delete profiles",
    "✅ Auto-remembers last used profile",
    "",
    "USAGE:",
    "1. Login once with credentials",
    "2. Check 'Remember credentials' checkbox",
    "3. Next time: Select profile from dropdown",
    "4. System auto-fills User ID & Password",
    "5. Just solve captcha and login!"
])

# SLIDE 12: Session History
add_content_slide(prs, "Session History Feature", [
    "📜 RECENT SESSIONS (Last 10)",
    "",
    "QUICK ACCESS DROPDOWN:",
    "├─ rr20251117_1430 (350 khatas, 2h 15m ago)",
    "├─ rr20251116_0945 (125 khatas, Yesterday)",
    "├─ rr20251115_1620 (89 khatas, 2 days ago)",
    "└─ ... (up to 10 sessions)",
    "",
    "FEATURES:",
    "✅ Quick folder access with one click",
    "✅ Shows: Date, Time, Khata count, Time ago",
    "✅ Auto-sorts by date (newest first)",
    "✅ Stores session metadata (villages processed, stats)",
    "",
    "KEYBOARD SHORTCUT:",
    "Press Ctrl+H to view session history dropdown"
])

# SLIDE 13: Troubleshooting
add_content_slide(prs, "Troubleshooting Common Issues", [
    "❓ COMMON ISSUES & SOLUTIONS",
    "",
    "Issue: Chrome won't launch",
    "  → Close all Chrome windows and try again",
    "",
    "Issue: Login fails / captcha error",
    "  → Solve captcha manually, system waits automatically",
    "",
    "Issue: Khata dropdown appears empty",
    "  → System auto-waits up to 5 seconds for dropdown to populate",
    "",
    "Issue: PDFs not saving",
    "  → Check session folder path, ensure sufficient disk space",
    "",
    "Issue: 'Already Paid' appears multiple times",
    "  → This is normal - system detected and skipped correctly",
    "",
    "Issue: Some khatas show error status",
    "  → Check error_log.txt for detailed error messages",
    "  → System automatically retries failed khatas 3 times",
    "",
    "⚠️ IF PORTAL CRASHES:",
    "• Note the last successfully processed khata",
    "• Close and restart MOHARIAR",
    "• Uncheck already-processed khatas",
    "• Resume processing from next khata"
])

# SLIDE 14: Keyboard Shortcuts
add_content_slide(prs, "Keyboard Shortcuts ⌨️", [
    "SPEED UP YOUR WORKFLOW:",
    "",
    "Ctrl + S  →  Start / Engage Protocol",
    "Ctrl + E  →  Export Results to Excel",
    "Ctrl + O  →  Open Current Session Folder",
    "Ctrl + H  →  View Session History",
    "Ctrl + Q  →  Stop Processing (abort)",
    "Ctrl + M  →  Merge PDFs Manually",
    "Ctrl + ?  →  Open Help Guide (this presentation)",
    "",
    "NAVIGATION:",
    "Space     →  Toggle current checkbox selection",
    "Enter     →  Confirm dialog boxes",
    "Esc       →  Close dialog boxes",
    "",
    "TIPS:",
    "• Shortcuts work throughout the application",
    "• Case-insensitive (Ctrl+s or Ctrl+S both work)",
    "• Most dialogs can be closed with Esc key"
])

# SLIDE 15: Technical Architecture
add_content_slide(prs, "Technical Architecture 🔧", [
    "TECHNOLOGY STACK:",
    "├─ GUI: CustomTkinter 5.2.2 (Modern dark theme)",
    "├─ Automation: Selenium WebDriver 4.16.0",
    "├─ Smart Delays: WebDriverWait (adaptive timing)",
    "├─ PDF Generation: Chrome DevTools Protocol (CDP)",
    "├─ PDF Merging: PyPDF2 3.0.1",
    "├─ Excel Processing: Pandas 2.1.0 + OpenPyXL 3.1.2",
    "└─ Notifications: win10toast (Windows toast popups)",
    "",
    "KEY INNOVATIONS:",
    "✅ WebDriverWait - Waits only as long as needed (40-50% faster)",
    "✅ CDP printToPDF - Direct browser PDF generation (no print dialogs)",
    "✅ JavaScript Alert Handler - Detects alerts before page interaction",
    "✅ Smart Retry - 3 attempts with exponential backoff (5s, 15s, 30s)",
    "✅ Khata Sanitization - Handles special chars (/ → _)",
    "",
    "PERFORMANCE:",
    "• Old method: 21s per khata (fixed delays)",
    "• MOHARIAR: 8-12s per khata (adaptive delays)",
    "• PDF Success Rate: 95%+",
    "• Already-Paid Detection: 100% accurate"
])

# SLIDE 16: Thank You
add_title_slide(
    prs,
    "⚡ THANK YOU! ⚡",
    "MOHARIAR v2.0\n\n👨‍💻 Designed by: SUSHANT\n📅 November 2025\n\n🙏 Thank you for using MOHARIAR!\n\nQuestions? Feedback? Suggestions?\nFeel free to reach out!"
)

# Save presentation
prs.save('MOHARIAR_Guide.pptx')
print("SUCCESS: PowerPoint presentation created: MOHARIAR_Guide.pptx")
print("\nNext Steps:")
print("1. Take screenshots of MOHARIAR UI:")
print("   - Slide 5: Full application window (3-panel layout)")
print("   - Slide 10: Right panel showing live statistics")
print("2. Open MOHARIAR_Guide.pptx in PowerPoint")
print("3. Insert screenshots where indicated")
print("4. Review and adjust formatting as needed")
print("5. Export as PDF: File -> Save As -> PDF")
